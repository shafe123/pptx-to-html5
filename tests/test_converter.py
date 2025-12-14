"""Test suite for PowerPoint to HTML5 converter."""

import tempfile
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_to_html5.converter import PowerPointToHTML5Converter


@pytest.fixture
def sample_pptx() -> Path:
    """Create a sample PowerPoint presentation for testing.

    Returns:
        Path to the temporary PowerPoint file
    """
    prs = Presentation()
    prs.slide_width = 9144000  # 10 inches
    prs.slide_height = 6858000  # 7.5 inches

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Sample Subtitle"

    # Add a content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Content Slide"
    tf = body_shape.text_frame
    tf.text = "First bullet point"

    # Save to temporary file
    with tempfile.NamedTemporaryFile(
        suffix=".pptx", delete=False
    ) as tmp_file:
        tmp_path = Path(tmp_file.name)
        prs.save(tmp_file.name)

    yield tmp_path

    # Cleanup
    if tmp_path.exists():
        tmp_path.unlink()


@pytest.fixture
def output_dir() -> Path:
    """Create a temporary output directory.

    Returns:
        Path to the temporary directory
    """
    with tempfile.TemporaryDirectory() as tmp_dir:
        yield Path(tmp_dir)


class TestPowerPointToHTML5Converter:
    """Test cases for PowerPointToHTML5Converter class."""

    def test_init_valid_file(self, sample_pptx: Path) -> None:
        """Test initialization with a valid PowerPoint file."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        assert converter.pptx_path == sample_pptx
        assert converter.presentation is not None

    def test_init_file_not_found(self) -> None:
        """Test initialization with a non-existent file."""
        with pytest.raises(FileNotFoundError):
            PowerPointToHTML5Converter("nonexistent.pptx")

    def test_init_invalid_extension(self, tmp_path: Path) -> None:
        """Test initialization with an invalid file extension."""
        invalid_file = tmp_path / "test.txt"
        invalid_file.touch()
        with pytest.raises(ValueError, match="must be a .pptx file"):
            PowerPointToHTML5Converter(invalid_file)

    def test_init_invalid_pptx(self, tmp_path: Path) -> None:
        """Test initialization with an invalid PowerPoint file."""
        invalid_pptx = tmp_path / "invalid.pptx"
        invalid_pptx.write_text("not a valid pptx file")
        with pytest.raises(ValueError, match="Invalid PowerPoint file"):
            PowerPointToHTML5Converter(invalid_pptx)

    def test_extract_slide_content(self, sample_pptx: Path) -> None:
        """Test extracting content from a slide."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        slides = list(converter.presentation.slides)
        assert len(slides) >= 1

        content = converter._extract_slide_content(slides[0])
        assert isinstance(content, dict)
        assert "title" in content
        assert "shapes" in content
        assert "notes" in content
        assert "slide_width" in content
        assert "slide_height" in content
        # Check that shapes were extracted
        assert len(content["shapes"]) > 0
        # Verify at least one shape contains the expected text
        shape_texts = [s.get("text", "") for s in content["shapes"]]
        assert "Test Presentation" in shape_texts

    def test_slide_to_image(self, sample_pptx: Path) -> None:
        """Test converting a slide to an image."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        slides = list(converter.presentation.slides)
        image_data = converter._slide_to_image(slides[0], 0)

        assert isinstance(image_data, str)
        assert image_data.startswith("data:image/png;base64,")
        assert len(image_data) > 100  # Should have substantial content

    def test_convert_creates_files(
        self, sample_pptx: Path, output_dir: Path
    ) -> None:
        """Test that convert creates all necessary files."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        html_path = converter.convert(output_dir)

        assert html_path.exists()
        assert html_path.name == "index.html"
        assert (output_dir / "styles.css").exists()
        assert (output_dir / "script.js").exists()

    def test_convert_html_content(
        self, sample_pptx: Path, output_dir: Path
    ) -> None:
        """Test that the generated HTML contains expected content."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        html_path = converter.convert(output_dir)

        html_content = html_path.read_text()
        assert "<!DOCTYPE html>" in html_content
        assert "Test Presentation" in html_content
        assert "slide" in html_content.lower()
        assert "styles.css" in html_content
        assert "script.js" in html_content

    def test_convert_with_notes(
        self, sample_pptx: Path, output_dir: Path
    ) -> None:
        """Test conversion with speaker notes included."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        html_path = converter.convert(output_dir, include_notes=True)

        html_content = html_path.read_text()
        assert html_content is not None

    def test_convert_without_notes(
        self, sample_pptx: Path, output_dir: Path
    ) -> None:
        """Test conversion without speaker notes."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        html_path = converter.convert(output_dir, include_notes=False)

        html_content = html_path.read_text()
        assert "speaker notes" not in html_content.lower()

    def test_generate_css(self, sample_pptx: Path) -> None:
        """Test CSS generation."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        css = converter._generate_css()

        assert isinstance(css, str)
        assert len(css) > 0
        assert ".slide" in css
        assert ".nav-btn" in css
        assert ".progress-bar" in css

    def test_generate_js(self, sample_pptx: Path) -> None:
        """Test JavaScript generation."""
        converter = PowerPointToHTML5Converter(sample_pptx)
        js = converter._generate_js()

        assert isinstance(js, str)
        assert len(js) > 0
        assert "currentSlide" in js
        assert "nextSlide" in js
        assert "previousSlide" in js
        assert "addEventListener" in js

    def test_convert_creates_output_directory(
        self, sample_pptx: Path, tmp_path: Path
    ) -> None:
        """Test that convert creates the output directory if it doesn't exist."""
        output_dir = tmp_path / "new_output"
        assert not output_dir.exists()

        converter = PowerPointToHTML5Converter(sample_pptx)
        html_path = converter.convert(output_dir)

        assert output_dir.exists()
        assert html_path.exists()

    def test_multiple_slides(self, tmp_path: Path, output_dir: Path) -> None:
        """Test conversion with multiple slides."""
        # Create a presentation with multiple slides
        prs = Presentation()
        for i in range(5):
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = f"Slide {i + 1}"

        pptx_path = tmp_path / "multi_slide.pptx"
        prs.save(str(pptx_path))

        converter = PowerPointToHTML5Converter(pptx_path)
        html_path = converter.convert(output_dir)

        html_content = html_path.read_text()
        # Check that all slides are present
        for i in range(5):
            assert f"Slide {i + 1}" in html_content
