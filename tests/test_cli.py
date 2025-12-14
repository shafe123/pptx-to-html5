"""Test suite for CLI module."""

import subprocess
import sys
import tempfile
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation

from pptx_to_html5.cli import main


@pytest.fixture
def sample_pptx() -> Path:
    """Create a sample PowerPoint presentation for testing.

    Returns:
        Path to the temporary PowerPoint file
    """
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "CLI Test Presentation"

    with tempfile.NamedTemporaryFile(
        suffix=".pptx", delete=False
    ) as tmp_file:
        tmp_path = Path(tmp_file.name)
        prs.save(tmp_file.name)

    yield tmp_path

    # Cleanup
    if tmp_path.exists():
        tmp_path.unlink()


class TestCLI:
    """Test cases for CLI functionality."""

    def test_main_success(
        self, sample_pptx: Path, tmp_path: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        """Test successful conversion via CLI."""
        output_dir = tmp_path / "output"

        with patch.object(
            sys,
            "argv",
            ["pptx-to-html", str(sample_pptx), "-o", str(output_dir)],
        ):
            exit_code = main()

        assert exit_code == 0
        captured = capsys.readouterr()
        assert "Successfully converted" in captured.out
        assert str(output_dir / "index.html") in captured.out

    def test_main_file_not_found(self, capsys: pytest.CaptureFixture[str]) -> None:
        """Test CLI with non-existent file."""
        with patch.object(
            sys, "argv", ["pptx-to-html", "nonexistent.pptx"]
        ):
            exit_code = main()

        assert exit_code == 1
        captured = capsys.readouterr()
        assert "Error" in captured.err

    def test_main_with_notes(
        self, sample_pptx: Path, tmp_path: Path
    ) -> None:
        """Test CLI with --include-notes flag."""
        output_dir = tmp_path / "output"

        with patch.object(
            sys,
            "argv",
            [
                "pptx-to-html",
                str(sample_pptx),
                "-o",
                str(output_dir),
                "--include-notes",
            ],
        ):
            exit_code = main()

        assert exit_code == 0
        assert (output_dir / "index.html").exists()

    def test_main_default_output(
        self, sample_pptx: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        """Test CLI with default output directory."""
        with patch.object(sys, "argv", ["pptx-to-html", str(sample_pptx)]):
            exit_code = main()

        assert exit_code == 0
        captured = capsys.readouterr()
        assert "Successfully converted" in captured.out

    def test_main_invalid_file(
        self, tmp_path: Path, capsys: pytest.CaptureFixture[str]
    ) -> None:
        """Test CLI with invalid PowerPoint file."""
        invalid_file = tmp_path / "invalid.pptx"
        invalid_file.write_text("not a valid pptx")

        with patch.object(sys, "argv", ["pptx-to-html", str(invalid_file)]):
            exit_code = main()

        assert exit_code == 1
        captured = capsys.readouterr()
        assert "Error" in captured.err

    def test_help_output(self) -> None:
        """Test that --help displays usage information."""
        result = subprocess.run(
            [sys.executable, "-m", "pptx_to_html5.cli", "--help"],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "Convert PowerPoint presentations to HTML5" in result.stdout
        assert "input" in result.stdout
        assert "--output" in result.stdout

    def test_version_output(self) -> None:
        """Test that --version displays version information."""
        result = subprocess.run(
            [sys.executable, "-m", "pptx_to_html5.cli", "--version"],
            capture_output=True,
            text=True,
        )

        assert result.returncode == 0
        assert "0.1.0" in result.stdout
