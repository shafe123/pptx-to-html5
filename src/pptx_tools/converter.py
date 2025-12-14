"""Core PowerPoint to HTML5 converter."""

import base64
import io
from pathlib import Path
from typing import Any

from jinja2 import Template
from PIL import Image
from pptx import Presentation
from pptx.slide import Slide


class PowerPointToHTML5Converter:
    """Convert PowerPoint presentations to HTML5 websites."""

    def __init__(self, pptx_path: str | Path) -> None:
        """Initialize the converter with a PowerPoint file.

        Args:
            pptx_path: Path to the PowerPoint file (.pptx)

        Raises:
            FileNotFoundError: If the PowerPoint file doesn't exist
            ValueError: If the file is not a valid PowerPoint file
        """
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"File not found: {pptx_path}")
        if not self.pptx_path.suffix.lower() == ".pptx":
            raise ValueError(f"File must be a .pptx file: {pptx_path}")

        try:
            self.presentation = Presentation(str(self.pptx_path))
        except Exception as e:
            raise ValueError(f"Invalid PowerPoint file: {e}") from e

    def _slide_to_image(self, slide: Slide, slide_number: int) -> str:
        """Convert a slide to a base64-encoded PNG image.

        Args:
            slide: The slide to convert
            slide_number: The slide number (for identification)

        Returns:
            Base64-encoded PNG image as a data URI
        """
        # Create an image from the slide
        # Since python-pptx doesn't provide direct rendering, we'll export metadata
        # and use a placeholder approach. For production, you'd use a library like
        # aspose.slides or convert via LibreOffice/unoconv
        img = Image.new("RGB", (1280, 720), color=(255, 255, 255))

        # Convert to base64
        buffer = io.BytesIO()
        img.save(buffer, format="PNG")
        img_str = base64.b64encode(buffer.getvalue()).decode()
        return f"data:image/png;base64,{img_str}"

    def _extract_slide_content(self, slide: Slide) -> dict[str, Any]:
        """Extract text and content from a slide.

        Args:
            slide: The slide to extract content from

        Returns:
            Dictionary containing slide content
        """
        content: dict[str, Any] = {
            "title": "",
            "text": [],
            "notes": "",
        }

        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    # First text element is often the title
                    if not content["title"] and hasattr(shape, "text_frame"):
                        content["title"] = text
                    else:
                        content["text"].append(text)

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                content["notes"] = notes_slide.notes_text_frame.text.strip()

        return content

    def convert(
        self, output_dir: str | Path, include_notes: bool = False
    ) -> Path:
        """Convert the PowerPoint presentation to an HTML5 website.

        Args:
            output_dir: Directory where the HTML5 website will be created
            include_notes: Whether to include speaker notes in the output

        Returns:
            Path to the generated index.html file

        Raises:
            OSError: If the output directory cannot be created
        """
        output_path = Path(output_dir)
        output_path.mkdir(parents=True, exist_ok=True)

        # Extract slides
        slides_data = []
        for i, slide in enumerate(self.presentation.slides):
            slide_content = self._extract_slide_content(slide)
            slide_data = {
                "number": i + 1,
                "image": self._slide_to_image(slide, i),
                "title": slide_content["title"],
                "text": slide_content["text"],
                "notes": slide_content["notes"] if include_notes else "",
            }
            slides_data.append(slide_data)

        # Generate HTML
        html_content = self._generate_html(slides_data, include_notes)
        html_path = output_path / "index.html"
        html_path.write_text(html_content, encoding="utf-8")

        # Generate CSS
        css_content = self._generate_css()
        css_path = output_path / "styles.css"
        css_path.write_text(css_content, encoding="utf-8")

        # Generate JavaScript
        js_content = self._generate_js()
        js_path = output_path / "script.js"
        js_path.write_text(js_content, encoding="utf-8")

        return html_path

    def _generate_html(
        self, slides_data: list[dict[str, Any]], include_notes: bool
    ) -> str:
        """Generate HTML content for the presentation.

        Args:
            slides_data: List of slide data dictionaries
            include_notes: Whether to include speaker notes

        Returns:
            HTML content as a string
        """
        template = Template(
            """<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <link rel="stylesheet" href="styles.css">
</head>
<body>
    <div class="presentation-container">
        <div class="slides-wrapper">
            {% for slide in slides %}
            <div class="slide{% if loop.first %} active{% endif %}" data-slide="{{ slide.number }}">
                <div class="slide-content">
                    {% if slide.title %}
                    <h1 class="slide-title">{{ slide.title }}</h1>
                    {% endif %}
                    {% if slide.image %}
                    <div class="slide-image">
                        <img src="{{ slide.image }}" alt="Slide {{ slide.number }}">
                    </div>
                    {% endif %}
                    {% if slide.text %}
                    <div class="slide-text">
                        {% for text in slide.text %}
                        <p>{{ text }}</p>
                        {% endfor %}
                    </div>
                    {% endif %}
                    {% if include_notes and slide.notes %}
                    <div class="slide-notes">
                        <h3>Speaker Notes:</h3>
                        <p>{{ slide.notes }}</p>
                    </div>
                    {% endif %}
                </div>
            </div>
            {% endfor %}
        </div>

        <div class="controls">
            <button id="prevBtn" class="nav-btn">← Previous</button>
            <span class="slide-counter">
                <span id="currentSlide">1</span> / <span id="totalSlides">{{ slides|length }}</span>
            </span>
            <button id="nextBtn" class="nav-btn">Next →</button>
        </div>

        <div class="progress-bar">
            <div class="progress-fill" id="progressFill"></div>
        </div>
    </div>

    <script src="script.js"></script>
</body>
</html>"""
        )

        title = (
            slides_data[0]["title"]
            if slides_data and slides_data[0]["title"]
            else "Presentation"
        )

        return template.render(
            title=title, slides=slides_data, include_notes=include_notes
        )

    def _generate_css(self) -> str:
        """Generate CSS styles for the presentation.

        Returns:
            CSS content as a string
        """
        return """* {
    margin: 0;
    padding: 0;
    box-sizing: border-box;
}

body {
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, Cantarell, sans-serif;
    background: #1a1a1a;
    color: #333;
    overflow: hidden;
}

.presentation-container {
    width: 100vw;
    height: 100vh;
    display: flex;
    flex-direction: column;
    justify-content: center;
    align-items: center;
}

.slides-wrapper {
    flex: 1;
    width: 100%;
    position: relative;
    display: flex;
    align-items: center;
    justify-content: center;
    padding: 2rem;
}

.slide {
    display: none;
    width: 100%;
    max-width: 1200px;
    height: 100%;
    background: white;
    border-radius: 8px;
    box-shadow: 0 10px 40px rgba(0, 0, 0, 0.3);
    padding: 3rem;
    overflow-y: auto;
    animation: fadeIn 0.3s ease-in;
}

.slide.active {
    display: block;
}

@keyframes fadeIn {
    from {
        opacity: 0;
        transform: translateY(20px);
    }
    to {
        opacity: 1;
        transform: translateY(0);
    }
}

.slide-content {
    max-width: 900px;
    margin: 0 auto;
}

.slide-title {
    font-size: 2.5rem;
    color: #2c3e50;
    margin-bottom: 2rem;
    border-bottom: 3px solid #3498db;
    padding-bottom: 0.5rem;
}

.slide-image {
    margin: 2rem 0;
    text-align: center;
}

.slide-image img {
    max-width: 100%;
    height: auto;
    border-radius: 4px;
    box-shadow: 0 2px 8px rgba(0, 0, 0, 0.1);
}

.slide-text {
    font-size: 1.2rem;
    line-height: 1.8;
    color: #555;
}

.slide-text p {
    margin-bottom: 1rem;
}

.slide-notes {
    margin-top: 2rem;
    padding: 1.5rem;
    background: #f8f9fa;
    border-left: 4px solid #3498db;
    border-radius: 4px;
}

.slide-notes h3 {
    color: #2c3e50;
    margin-bottom: 0.5rem;
    font-size: 1.2rem;
}

.slide-notes p {
    color: #666;
    font-size: 1rem;
    line-height: 1.6;
}

.controls {
    display: flex;
    align-items: center;
    gap: 2rem;
    padding: 1.5rem;
    background: rgba(255, 255, 255, 0.95);
    border-radius: 50px;
    box-shadow: 0 4px 20px rgba(0, 0, 0, 0.2);
    margin-bottom: 2rem;
}

.nav-btn {
    padding: 0.8rem 1.5rem;
    font-size: 1rem;
    background: #3498db;
    color: white;
    border: none;
    border-radius: 25px;
    cursor: pointer;
    transition: all 0.3s ease;
    font-weight: 500;
}

.nav-btn:hover {
    background: #2980b9;
    transform: translateY(-2px);
    box-shadow: 0 4px 12px rgba(52, 152, 219, 0.3);
}

.nav-btn:active {
    transform: translateY(0);
}

.nav-btn:disabled {
    background: #bdc3c7;
    cursor: not-allowed;
    transform: none;
}

.slide-counter {
    font-size: 1.1rem;
    color: #555;
    font-weight: 500;
    min-width: 80px;
    text-align: center;
}

.progress-bar {
    width: 100%;
    height: 4px;
    background: rgba(255, 255, 255, 0.3);
    position: fixed;
    bottom: 0;
    left: 0;
}

.progress-fill {
    height: 100%;
    background: #3498db;
    transition: width 0.3s ease;
}

/* Keyboard shortcuts hint */
body::before {
    content: "Use ← → keys or click buttons to navigate";
    position: fixed;
    top: 1rem;
    right: 1rem;
    padding: 0.5rem 1rem;
    background: rgba(255, 255, 255, 0.9);
    border-radius: 4px;
    font-size: 0.85rem;
    color: #666;
    z-index: 1000;
}

/* Responsive design */
@media (max-width: 768px) {
    .slide {
        padding: 1.5rem;
    }

    .slide-title {
        font-size: 1.8rem;
    }

    .slide-text {
        font-size: 1rem;
    }

    .controls {
        gap: 1rem;
        padding: 1rem;
    }

    .nav-btn {
        padding: 0.6rem 1rem;
        font-size: 0.9rem;
    }

    body::before {
        font-size: 0.75rem;
        padding: 0.4rem 0.8rem;
    }
}
"""

    def _generate_js(self) -> str:
        """Generate JavaScript for slide navigation and interactivity.

        Returns:
            JavaScript content as a string
        """
        return """// Presentation navigation
let currentSlide = 1;
let totalSlides = 0;

// Initialize
document.addEventListener('DOMContentLoaded', function() {
    const slides = document.querySelectorAll('.slide');
    totalSlides = slides.length;

    document.getElementById('totalSlides').textContent = totalSlides;
    updateSlide();

    // Navigation buttons
    document.getElementById('prevBtn').addEventListener('click', previousSlide);
    document.getElementById('nextBtn').addEventListener('click', nextSlide);

    // Keyboard navigation
    document.addEventListener('keydown', function(e) {
        if (e.key === 'ArrowLeft' || e.key === 'ArrowUp') {
            previousSlide();
        } else if (e.key === 'ArrowRight' || e.key === 'ArrowDown' || e.key === ' ') {
            e.preventDefault();
            nextSlide();
        } else if (e.key === 'Home') {
            goToSlide(1);
        } else if (e.key === 'End') {
            goToSlide(totalSlides);
        }
    });

    // Touch/swipe support
    let touchStartX = 0;
    let touchEndX = 0;

    document.addEventListener('touchstart', function(e) {
        touchStartX = e.changedTouches[0].screenX;
    });

    document.addEventListener('touchend', function(e) {
        touchEndX = e.changedTouches[0].screenX;
        handleSwipe();
    });

    function handleSwipe() {
        const swipeThreshold = 50;
        const diff = touchStartX - touchEndX;

        if (Math.abs(diff) > swipeThreshold) {
            if (diff > 0) {
                nextSlide();
            } else {
                previousSlide();
            }
        }
    }
});

function updateSlide() {
    const slides = document.querySelectorAll('.slide');

    // Hide all slides
    slides.forEach(slide => slide.classList.remove('active'));

    // Show current slide
    slides[currentSlide - 1].classList.add('active');

    // Update counter
    document.getElementById('currentSlide').textContent = currentSlide;

    // Update progress bar
    const progress = (currentSlide / totalSlides) * 100;
    document.getElementById('progressFill').style.width = progress + '%';

    // Update button states
    document.getElementById('prevBtn').disabled = currentSlide === 1;
    document.getElementById('nextBtn').disabled = currentSlide === totalSlides;
}

function nextSlide() {
    if (currentSlide < totalSlides) {
        currentSlide++;
        updateSlide();
    }
}

function previousSlide() {
    if (currentSlide > 1) {
        currentSlide--;
        updateSlide();
    }
}

function goToSlide(slideNumber) {
    if (slideNumber >= 1 && slideNumber <= totalSlides) {
        currentSlide = slideNumber;
        updateSlide();
    }
}

// Prevent default space bar scrolling
window.addEventListener('keydown', function(e) {
    if (e.key === ' ' && e.target === document.body) {
        e.preventDefault();
    }
});
"""
